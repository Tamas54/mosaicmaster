# external_converter_helper.py
"""
Külső dokumentum konverter segédmodul problémás formátumokhoz
DOC, PPT, PPTX, MOBI támogatás graceful fallback-ekkel
+ GPT-4o Vision OCR támogatás képekhez és PDF-ekhez

TELJES IMPLEMENTÁCIÓ - a működő MOBI parser alapján + Vision OCR
Railway-kompatibilis, minimális dependencies
"""

import asyncio
import logging
import tempfile
import shutil
import subprocess
import sys
import struct
import base64
import os
from pathlib import Path
from typing import Dict, Any, Optional, List
import aiofiles
from openai import OpenAI
from PIL import Image
from dotenv import load_dotenv

logger = logging.getLogger(__name__)

# Feature flags - ezek jelzik, hogy mely funkciók érhetők el
HAS_DOC_SUPPORT = False
HAS_PPTX_SUPPORT = False
HAS_PPT_SUPPORT = False
HAS_MOBI_SUPPORT = False
HAS_BEAUTIFULSOUP = False
HAS_LIBREOFFICE = False
HAS_VISION_OCR = False

# OCR beállítások
MODEL_NAME = "gpt-4o-mini"          # Vision-képes modell
DETAIL_LEVEL = "high"               # "low" olcsóbb, de pontatlanabb lehet
MAX_TOKENS = 8000                   # OCR-válasz hossz - 1-2 oldalnyi szöveghez

# Támogatott képformátumok OCR-hez
SUPPORTED_IMAGE_EXTENSIONS = {'.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.webp', '.heic'}

# MOBI támogatási módszerek
MOBI_METHODS = []

# Package verzió információk
PACKAGE_VERSIONS = {}

# Conditional imports - ha hibáznak, az adott feature letiltódik
try:
    import docx2txt
    HAS_DOC_SUPPORT = True
    PACKAGE_VERSIONS["docx2txt"] = getattr(docx2txt, '__version__', 'unknown')
    logger.info("DOC support available via docx2txt")
except ImportError:
    logger.warning("docx2txt not available - DOC support disabled")

try:
    from pptx import Presentation
    import pptx
    HAS_PPTX_SUPPORT = True
    PACKAGE_VERSIONS["python-pptx"] = getattr(pptx, '__version__', 'unknown')
    logger.info("PPTX support available via python-pptx")
except ImportError:
    logger.warning("python-pptx not available - PPTX support disabled")

try:
    import olefile
    HAS_PPT_SUPPORT = True
    PACKAGE_VERSIONS["olefile"] = getattr(olefile, '__version__', 'unknown')
    logger.info("PPT support available via olefile")
except ImportError:
    logger.warning("olefile not available - PPT support disabled")

try:
    from bs4 import BeautifulSoup
    HAS_BEAUTIFULSOUP = True
    import bs4
    PACKAGE_VERSIONS["beautifulsoup4"] = getattr(bs4, '__version__', 'unknown')
    logger.info("BeautifulSoup available for HTML parsing")
except ImportError:
    logger.warning("BeautifulSoup not available - HTML parsing limited")

# MOBI támogatás - a működő parser alapján
try:
    import mobi
    MOBI_METHODS.append("mobi")
    PACKAGE_VERSIONS["mobi"] = getattr(mobi, '__version__', 'unknown')
    logger.info("MOBI support via 'mobi' library")
except ImportError:
    pass

try:
    from ebooklib import epub
    MOBI_METHODS.append("ebooklib")
    import ebooklib
    PACKAGE_VERSIONS["ebooklib"] = getattr(ebooklib, '__version__', 'unknown')
    logger.info("EPUB support via 'ebooklib' (limited MOBI)")
except ImportError:
    pass

HAS_MOBI_SUPPORT = len(MOBI_METHODS) > 0

# OpenAI Vision OCR support
try:
    load_dotenv()
    if os.getenv("OPENAI_API_KEY"):
        openai_client = OpenAI()
        HAS_VISION_OCR = True
        logger.info("OpenAI Vision OCR support available")
    else:
        logger.warning("OpenAI Vision OCR disabled - missing OPENAI_API_KEY")
        openai_client = None
except Exception as e:
    logger.warning(f"OpenAI Vision OCR disabled - {str(e)}")
    openai_client = None

# LibreOffice support disabled for Railway deployment
HAS_LIBREOFFICE = False
logger.info("LibreOffice support disabled for Railway deployment")


class ExternalConverterHelper:
    """Külső konverter segédosztály problémás formátumokhoz"""
    
    @staticmethod
    def _image_to_base64(image_path: Path) -> str:
        """
        Betölt egy képet, biztosítja a JPEG/PNG formátumot, és base64-re kódolja.
        """
        # Ha nem JPEG/PNG, konvertáljuk (pl. HEIC → JPEG) a PIL segítségével
        if image_path.suffix.lower() not in {".jpg", ".jpeg", ".png"}:
            img = Image.open(image_path)
            temp_path = image_path.with_suffix(".jpg")
            img.save(temp_path, format="JPEG", quality=95)
            image_path = temp_path

        with open(image_path, "rb") as f:
            return base64.b64encode(f.read()).decode("utf-8")
    
    @staticmethod
    async def ocr_image_with_vision(input_path: Path, output_path: Path, language_hint: str = "any language") -> Dict[str, Any]:
        """Kép OCR feldolgozása GPT-4o Vision segítségével"""
        if not HAS_VISION_OCR:
            raise RuntimeError("Vision OCR support not available - missing OpenAI API key")
        
        try:
            loop = asyncio.get_event_loop()
            
            def process_with_vision():
                try:
                    logger.debug(f"Processing image with Vision OCR: {input_path}")
                    
                    # Kép base64 kódolása
                    b64_image = ExternalConverterHelper._image_to_base64(input_path)
                    
                    # OCR kérés Vision API-hoz - multilingual támogatással
                    if language_hint == "any language":
                        prompt_text = "Please read and transcribe ALL text visible in this image exactly as you see it. The text may be in ANY language (including but not limited to: English, Hungarian, German, French, Spanish, Italian, Russian, Chinese, Japanese, Korean, Arabic, Hindi, Portuguese, Dutch, Swedish, Norwegian, Danish, Finnish, Polish, Czech, Turkish, Greek, Hebrew, Thai, Vietnamese, Indonesian, and many others). Preserve the original formatting, line breaks, and layout as much as possible. If there's no text in the image, describe what you see."
                    else:
                        prompt_text = f"Please read and transcribe ALL text visible in this image exactly as you see it. The text is expected to be in {language_hint}. Preserve the original formatting and layout. If there's no text in the image, describe what you see."
                    
                    response = openai_client.chat.completions.create(
                        model=MODEL_NAME,
                        max_tokens=MAX_TOKENS,
                        messages=[
                            {
                                "role": "user",
                                "content": [
                                    {"type": "text", "text": prompt_text},
                                    {
                                        "type": "image_url",
                                        "image_url": {
                                            "url": f"data:image/jpeg;base64,{b64_image}",
                                            "detail": DETAIL_LEVEL
                                        },
                                    },
                                ],
                            }
                        ],
                    )
                    
                    # Eredmény kivonása
                    ocr_text = response.choices[0].message.content.strip()
                    
                    # Token használat logolása
                    if hasattr(response, 'usage') and response.usage:
                        usage = response.usage
                        logger.info(f"Vision OCR token usage: {usage.prompt_tokens} prompt + {usage.completion_tokens} completion = {usage.total_tokens} total")
                    
                    logger.info(f"Vision OCR extraction successful - {len(ocr_text)} characters")
                    return ocr_text
                    
                except Exception as e:
                    logger.error(f"Vision OCR processing failed: {str(e)}")
                    raise Exception(f"Vision OCR processing failed: {str(e)}")
            
            # Aszinkron feldolgozás
            text = await loop.run_in_executor(None, process_with_vision)
            
            # TXT fájl mentése
            async with aiofiles.open(output_path, 'w', encoding='utf-8') as f:
                await f.write(text)
                
            return {"converted": True, "method": "vision_ocr", "characters": len(text)}
            
        except Exception as e:
            logger.error(f"Vision OCR conversion error: {str(e)}")
            raise Exception(f"Vision OCR conversion failed: {str(e)}")

    @staticmethod
    async def convert_doc_to_txt(input_path: Path, output_path: Path) -> Dict[str, Any]:
        """DOC → TXT konverzió docx2txt vagy olefile használatával"""
        if not HAS_DOC_SUPPORT:
            raise RuntimeError("DOC support not available - docx2txt package missing")
        
        try:
            loop = asyncio.get_event_loop()
            
            def extract_doc_text():
                try:
                    logger.debug(f"Attempting DOC extraction with docx2txt: {input_path}")
                    text = docx2txt.process(str(input_path))
                    
                    if not text or len(text.strip()) < 10:
                        raise Exception("Empty or invalid DOC content extracted")
                    
                    # Szöveg tisztítása és normalizálása
                    import re
                    text = re.sub(r'\n\s*\n', '\n\n', text)  # Dupla sortörések megtartása
                    text = re.sub(r'[ \t]+', ' ', text)      # Szóközök normalizálása
                    text = text.strip()
                    
                    logger.info(f"DOC extraction successful - {len(text)} characters")
                    return text
                    
                except Exception as e:
                    logger.error(f"docx2txt failed: {str(e)}")
                    
                    # LibreOffice fallback disabled for Railway deployment
                    
                    raise Exception(f"DOC processing failed: {str(e)}")
                    
            text = await loop.run_in_executor(None, extract_doc_text)
            
            # TXT fájl mentése
            async with aiofiles.open(output_path, 'w', encoding='utf-8') as f:
                await f.write(text)
                
            return {"converted": True, "method": "docx2txt", "characters": len(text)}
            
        except Exception as e:
            # LibreOffice fallback disabled for Railway deployment
            
            logger.error(f"DOC conversion error: {str(e)}")
            raise Exception(f"DOC conversion failed: {str(e)}")

    @staticmethod
    async def convert_pptx_to_txt(input_path: Path, output_path: Path) -> Dict[str, Any]:
        """PPTX → TXT konverzió python-pptx használatával"""
        if not HAS_PPTX_SUPPORT:
            raise RuntimeError("PPTX support not available - python-pptx package missing")
        
        try:
            loop = asyncio.get_event_loop()
            
            def extract_pptx_text():
                try:
                    logger.debug(f"Attempting PPTX extraction: {input_path}")
                    prs = Presentation(str(input_path))
                    text_runs = []
                    slide_count = len(prs.slides)
                    
                    logger.debug(f"PPTX contains {slide_count} slides")
                    
                    for slide_num, slide in enumerate(prs.slides, 1):
                        text_runs.append(f"\n=== Slide {slide_num} / {slide_count} ===\n")
                        
                        slide_text_found = False
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text.strip():
                                text_runs.append(shape.text.strip())
                                slide_text_found = True
                        
                        # Ha nincs szöveg a slide-on
                        if not slide_text_found:
                            text_runs.append("[No text content on this slide]")
                    
                    full_text = "\n".join(text_runs)
                    
                    # Szöveg minőség ellenőrzése
                    meaningful_text = full_text.replace("=== Slide", "").replace("[No text content", "").strip()
                    if len(meaningful_text) < 20:
                        raise Exception("No meaningful text content found in PPTX")
                    
                    logger.info(f"PPTX extraction successful - {slide_count} slides, {len(full_text)} characters")
                    return full_text
                    
                except Exception as e:
                    logger.error(f"python-pptx failed: {str(e)}")
                    
                    # LibreOffice fallback disabled for Railway deployment
                    
                    raise Exception(f"PPTX text extraction failed: {str(e)}")
                    
            text = await loop.run_in_executor(None, extract_pptx_text)
            
            # TXT fájl mentése
            async with aiofiles.open(output_path, 'w', encoding='utf-8') as f:
                await f.write(text)
                
            return {"converted": True, "method": "python-pptx", "characters": len(text)}
            
        except Exception as e:
            # LibreOffice fallback disabled for Railway deployment
            
            logger.error(f"PPTX conversion error: {str(e)}")
            raise Exception(f"PPTX conversion failed: {str(e)}")

    @staticmethod
    async def convert_ppt_to_txt(input_path: Path, output_path: Path) -> Dict[str, Any]:
        """PPT → TXT konverzió (elsősorban LibreOffice-szal)"""
        # LibreOffice support disabled for Railway deployment
        logger.info("LibreOffice PPT conversion disabled for Railway deployment")
        
        # Pure Python PPT support using olefile
        if not HAS_PPT_SUPPORT:
            raise RuntimeError("PPT support not available - olefile package missing")
        
        try:
            loop = asyncio.get_event_loop()
            
            def extract_ppt_text():
                try:
                    logger.debug(f"Attempting PPT extraction (experimental): {input_path}")
                    import olefile
                    
                    if not olefile.isOleFile(str(input_path)):
                        raise Exception("Not a valid OLE file (PPT)")
                    
                    with olefile.OleFileIO(str(input_path)) as ole:
                        streams = ole.get_streams()
                        text_content = []
                        
                        logger.debug(f"PPT file contains {len(streams)} streams")
                        
                        for stream in streams:
                            stream_name = '/'.join(stream).lower()
                            # PowerPoint specifikus streamek keresése
                            if any(keyword in stream_name for keyword in ['powerpoint', 'document', 'slide', 'text']):
                                try:
                                    data = ole.get_stream(stream).read()
                                    
                                    # Próbálkozás különböző encoding-okkal
                                    for encoding in ['utf-16le', 'utf-8', 'cp1252']:
                                        try:
                                            decoded = data.decode(encoding, errors='ignore')
                                            
                                            # Szűrés - csak valódi szöveges részek
                                            import re
                                            clean_text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]', '', decoded)
                                            clean_text = re.sub(r'[ \t]+', ' ', clean_text)
                                            clean_text = re.sub(r'\n\s*\n\s*\n', '\n\n', clean_text)
                                            clean_text = clean_text.strip()
                                            
                                            # Ha van értelmes szöveg
                                            if len(clean_text) > 10 and any(c.isalpha() for c in clean_text):
                                                text_content.append(clean_text)
                                                break
                                        except:
                                            continue
                                except:
                                    continue
                        
                        if not text_content:
                            raise Exception("No readable text found in PPT file")
                        
                        # Duplikátumok eltávolítása és összefűzés
                        unique_texts = []
                        for text in text_content:
                            if text not in unique_texts and len(text) > 10:
                                unique_texts.append(text)
                        
                        combined_text = "\n\n".join(unique_texts)
                        
                        logger.info(f"PPT extraction successful (experimental) - {len(combined_text)} characters")
                        return combined_text
                        
                except Exception as e:
                    logger.error(f"olefile PPT processing failed: {str(e)}")
                    raise Exception(f"PPT text extraction failed: {str(e)}")
                    
            text = await loop.run_in_executor(None, extract_ppt_text)
            
            # TXT fájl mentése
            async with aiofiles.open(output_path, 'w', encoding='utf-8') as f:
                await f.write(text)
                
            return {"converted": True, "method": "olefile", "note": "experimental", "characters": len(text)}
            
        except Exception as e:
            logger.error(f"PPT conversion error: {str(e)}")
            raise Exception(f"PPT conversion failed: {str(e)}")

    @staticmethod
    async def convert_mobi_to_txt(input_path: Path, output_path: Path) -> Dict[str, Any]:
        """MOBI → TXT konverzió - a működő parser alapján"""
        if not HAS_MOBI_SUPPORT:
            raise RuntimeError("MOBI support not available - mobi or ebooklib package missing")
        
        # Próbáljuk végig a rendelkezésre álló módszereket
        for method in MOBI_METHODS:
            try:
                logger.info(f"Trying MOBI conversion with method: {method}")
                result = await ExternalConverterHelper._try_mobi_method(
                    method, input_path, output_path
                )
                if result:
                    return result
            except Exception as e:
                logger.warning(f"MOBI method {method} failed: {str(e)}")
                continue
        
        # Ha minden módszer hibázik, manual parsing
        try:
            logger.info("Trying manual MOBI parsing as last resort")
            result = await ExternalConverterHelper._manual_mobi_parsing(input_path, output_path)
            if result:
                return result
        except Exception as e:
            logger.error(f"Manual MOBI parsing failed: {str(e)}")
        
        raise Exception("All MOBI conversion methods failed")

    @staticmethod
    async def _try_mobi_method(method: str, input_path: Path, output_path: Path) -> Optional[Dict[str, Any]]:
        """Egy specifikus MOBI módszer kipróbálása - a működő parser alapján"""
        loop = asyncio.get_event_loop()
        
        if method == "mobi":
            def extract_with_mobi():
                import mobi
                
                logger.debug(f"Attempting MOBI extraction: {input_path}")
                
                # Mobi extraction ideiglenes könyvtárba
                tempdir, filepath = mobi.extract(str(input_path))
                logger.debug(f"MOBI extracted to temporary directory: {tempdir}")
                
                text_parts = []
                
                # Keresünk kinyert szöveges fájlokat
                temp_path = Path(tempdir)
                extracted_files = list(temp_path.rglob("*"))
                logger.debug(f"Found {len(extracted_files)} extracted files")
                
                for file in extracted_files:
                    if file.suffix.lower() in ['.txt', '.html', '.htm', '.xhtml']:
                        try:
                            with open(file, 'r', encoding='utf-8', errors='ignore') as f:
                                content = f.read()
                                
                            if len(content.strip()) > 50:
                                # HTML/XHTML fájlok esetén szöveg kinyerése
                                if file.suffix.lower() in ['.html', '.htm', '.xhtml']:
                                    if HAS_BEAUTIFULSOUP:
                                        soup = BeautifulSoup(content, 'html.parser')
                                        clean_text = soup.get_text()
                                    else:
                                        # Egyszerű regex ha nincs BeautifulSoup
                                        import re
                                        clean_text = re.sub(r'<[^>]+>', '', content)
                                        clean_text = re.sub(r'&[a-zA-Z0-9#]+;', ' ', clean_text)
                                    
                                    if len(clean_text.strip()) > 50:
                                        text_parts.append(clean_text.strip())
                                else:
                                    text_parts.append(content.strip())
                                    
                        except Exception as file_error:
                            logger.debug(f"Could not process file {file}: {str(file_error)}")
                            continue
                
                # Cleanup temporary directory
                try:
                    shutil.rmtree(tempdir, ignore_errors=True)
                except Exception as cleanup_error:
                    logger.warning(f"Could not clean up temp directory {tempdir}: {str(cleanup_error)}")
                
                if not text_parts:
                    raise Exception("No readable text content found in MOBI file")
                
                # Összes szöveges rész egyesítése
                combined_text = '\n\n'.join(text_parts)
                
                # Szöveg normalizálása
                import re
                combined_text = re.sub(r'\n\s*\n\s*\n', '\n\n', combined_text)
                combined_text = combined_text.strip()
                
                logger.info(f"MOBI extraction successful - {len(text_parts)} text sections, {len(combined_text)} characters")
                return combined_text
            
            text = await loop.run_in_executor(None, extract_with_mobi)
            
        elif method == "ebooklib":
            def extract_with_ebooklib():
                from ebooklib import epub
                
                logger.debug(f"Attempting EPUB-as-MOBI extraction: {input_path}")
                
                # EbookLib MOBI-t közvetlenül nem támogatja, de megpróbálhatjuk
                try:
                    book = epub.read_epub(str(input_path))
                except:
                    # MOBI nem támogatott közvetlenül az ebooklib-ben
                    logger.debug("EbookLib: MOBI format not directly supported")
                    return None
                
                text_parts = []
                for item in book.get_items():
                    if item.get_type() == epub.EpubHtml:
                        try:
                            content = item.get_content().decode('utf-8')
                            # HTML tartalom szövegkinyerése
                            if HAS_BEAUTIFULSOUP:
                                soup = BeautifulSoup(content, 'html.parser')
                                text_parts.append(soup.get_text())
                            else:
                                # Egyszerű regex ha nincs BeautifulSoup
                                import re
                                clean_text = re.sub(r'<[^>]+>', '', content)
                                text_parts.append(clean_text)
                        except Exception:
                            continue
                
                if text_parts:
                    return '\n\n'.join(text_parts)
                
                return None
            
            text = await loop.run_in_executor(None, extract_with_ebooklib)
        else:
            return None
        
        if text and len(text.strip()) > 50:
            async with aiofiles.open(output_path, 'w', encoding='utf-8') as f:
                await f.write(text)
            
            return {"converted": True, "method": method, "characters": len(text)}
        
        return None

    @staticmethod
    async def _manual_mobi_parsing(input_path: Path, output_path: Path) -> Optional[Dict[str, Any]]:
        """Manuális MOBI szöveg kinyerés - a működő parser alapján"""
        loop = asyncio.get_event_loop()
        
        def manual_extract():
            logger.debug("Manual MOBI parsing (experimental)...")
            
            with open(input_path, 'rb') as f:
                data = f.read()
            
            # MOBI header keresése
            mobi_start = data.find(b'MOBI')
            if mobi_start == -1:
                return None
            
            # UTF-8 szöveges részek keresése
            try:
                decoded = data.decode('utf-8', errors='ignore')
                
                # HTML tagek eltávolítása
                import re
                no_html = re.sub(r'<[^>]+>', '', decoded)
                
                # Kontroll karakterek eltávolítása  
                clean_text = re.sub(r'[\x00-\x1f\x7f-\x9f]', ' ', no_html)
                
                # Jelentéketlen karakterek eltávolítása
                clean_text = re.sub(r'[^\w\s.,!?;:\-\'"()]', '', clean_text)
                
                # Többszörös szóközök eltávolítása
                clean_text = re.sub(r'\s+', ' ', clean_text).strip()
                
                # Ha van értelmes szöveg (legalább 100 karakter)
                if len(clean_text) > 100:
                    return clean_text
                
            except Exception:
                pass
            
            return None
        
        text = await loop.run_in_executor(None, manual_extract)
        
        if text and len(text.strip()) > 50:
            async with aiofiles.open(output_path, 'w', encoding='utf-8') as f:
                await f.write(text)
            
            return {"converted": True, "method": "manual", "characters": len(text)}
        
        return None

    # LibreOffice conversion method removed for Railway deployment

    @staticmethod
    def get_available_conversions() -> Dict[str, bool]:
        """Elérhető konverziók lekérdezése"""
        return {
            "doc": HAS_DOC_SUPPORT,
            "pptx": HAS_PPTX_SUPPORT,
            "ppt": HAS_PPT_SUPPORT,
            "mobi": HAS_MOBI_SUPPORT,
            "vision_ocr": HAS_VISION_OCR
        }

    @staticmethod
    def get_status_report() -> Dict[str, Any]:
        """Részletes státusz riport"""
        available = ExternalConverterHelper.get_available_conversions()
        
        missing_packages = {}
        recommendations = []
        
        if not HAS_DOC_SUPPORT:
            missing_packages["docx2txt"] = "pip install docx2txt"
            recommendations.append("DOC files: Install docx2txt")
            
        if not HAS_PPTX_SUPPORT:
            missing_packages["python-pptx"] = "pip install python-pptx"
            recommendations.append("PPTX files: Install python-pptx")
            
        if not HAS_PPT_SUPPORT:
            missing_packages["olefile"] = "pip install olefile"
            recommendations.append("Legacy PPT files: Install olefile for basic support")
            
        if not HAS_MOBI_SUPPORT:
            missing_packages["mobi"] = "pip install mobi"
            missing_packages["ebooklib"] = "pip install ebooklib"
            recommendations.append("MOBI files: Install 'mobi' or 'ebooklib' package")
            
        if not HAS_VISION_OCR:
            missing_packages["openai"] = "pip install openai"
            recommendations.append("Vision OCR for images: Set OPENAI_API_KEY environment variable")
            
        return {
            "available_features": available,
            "missing_packages": missing_packages,
            "recommendations": recommendations,
            "package_versions": PACKAGE_VERSIONS,
            "mobi_methods": MOBI_METHODS,
            "beautifulsoup_available": HAS_BEAUTIFULSOUP,
            "libreoffice_available": HAS_LIBREOFFICE,
            "vision_ocr_available": HAS_VISION_OCR
        }

    @staticmethod
    def get_user_friendly_status() -> str:
        """Felhasználóbarát státusz üzenet"""
        available = ExternalConverterHelper.get_available_conversions()
        supported_formats = [fmt.upper() for fmt, sup in available.items() if sup]
        
        if len(supported_formats) == 5:
            return "All advanced formats supported (DOC, PPTX, PPT, MOBI, VISION_OCR) - Railway compatible"
        elif len(supported_formats) > 0:
            return f"Advanced format support: {', '.join(supported_formats)} - Railway compatible"
        else:
            return "No advanced format support available - install required packages"


# Convenience functions for document_processor.py integration
async def try_convert_external(source_format: str, input_path: Path, output_path: Path, **kwargs) -> Optional[Dict[str, Any]]:
    """
    Megpróbálja konvertálni a fájlt külső helperrel
    Returns None ha nem támogatott, egyébként a konverzió eredménye
    """
    try:
        logger.debug(f"Attempting external conversion: {source_format} → TXT")
        
        if source_format == "doc" and HAS_DOC_SUPPORT:
            return await ExternalConverterHelper.convert_doc_to_txt(input_path, output_path)
        elif source_format == "pptx" and HAS_PPTX_SUPPORT:
            return await ExternalConverterHelper.convert_pptx_to_txt(input_path, output_path)
        elif source_format == "ppt" and HAS_PPT_SUPPORT:
            return await ExternalConverterHelper.convert_ppt_to_txt(input_path, output_path)
        elif source_format == "mobi" and HAS_MOBI_SUPPORT:
            return await ExternalConverterHelper.convert_mobi_to_txt(input_path, output_path)
        elif source_format in ["image", "pdf_image"] and HAS_VISION_OCR:
            # Képformátumok és PDF-ek Vision OCR-rel
            language_hint = kwargs.get("language_hint", "any language")
            return await ExternalConverterHelper.ocr_image_with_vision(input_path, output_path, language_hint)
        else:
            logger.debug(f"External conversion not available for {source_format}")
            return None
            
    except Exception as e:
        logger.warning(f"External conversion failed for {source_format}: {str(e)}")
        return None

async def try_ocr_image(input_path: Path, output_path: Path, language_hint: str = "any language") -> Optional[Dict[str, Any]]:
    """
    Képek OCR feldolgozása Vision API-val
    Returns None ha nem támogatott, egyébként a konverzió eredménye
    """
    try:
        # Ellenőrizzük, hogy képfájl-e
        if input_path.suffix.lower() not in SUPPORTED_IMAGE_EXTENSIONS:
            logger.debug(f"File extension {input_path.suffix} not supported for OCR")
            return None
            
        if not HAS_VISION_OCR:
            logger.debug("Vision OCR not available")
            return None
            
        return await ExternalConverterHelper.ocr_image_with_vision(input_path, output_path, language_hint)
        
    except Exception as e:
        logger.warning(f"OCR processing failed: {str(e)}")
        return None


def get_external_support_info() -> str:
    """User-friendly információ a külső támogatásról"""
    return ExternalConverterHelper.get_user_friendly_status()


# Debug és troubleshooting függvények
def print_diagnostic_info():
    """Diagnosztikai információk kiírása"""
    status = ExternalConverterHelper.get_status_report()
    
    print("🔧 EXTERNAL CONVERTER HELPER - DIAGNOSTIC INFO")
    print("=" * 60)
    print(f"Available features: {status['available_features']}")
    print(f"Package versions: {status['package_versions']}")
    print(f"MOBI methods: {status['mobi_methods']}")
    print(f"BeautifulSoup available: {status['beautifulsoup_available']}")
    print(f"LibreOffice available: {status['libreoffice_available']}")
    
    if status['missing_packages']:
        print("Missing packages:")
        for pkg, cmd in status['missing_packages'].items():
            print(f"  {pkg}: {cmd}")
    
    if status['recommendations']:
        print("Recommendations:")
        for rec in status['recommendations']:
            print(f"  - {rec}")


# Module initialization
logger.info(f"External Converter Helper loaded - {get_external_support_info()}")

# MOBI struktúra elemzés helper (a working parser alapján)
def analyze_mobi_structure(mobi_path: Path) -> Dict[str, Any]:
    """MOBI fájl struktúra elemzése (pure Python) - troubleshooting-hoz"""
    try:
        with open(mobi_path, 'rb') as f:
            data = f.read(1024)
            
            if len(data) < 78:
                return {"error": "Too short for MOBI file"}
            
            # MOBI magic number ellenőrzése
            mobi_header_start = data.find(b'MOBI')
            if mobi_header_start == -1:
                return {"error": "MOBI magic number not found"}
            
            info = {
                "has_mobi_header": True,
                "mobi_header_offset": mobi_header_start,
                "file_size": mobi_path.stat().st_size,
            }
            
            try:
                # MOBI header hossza
                header_len = struct.unpack('>I', data[mobi_header_start+4:mobi_header_start+8])[0]
                info["header_length"] = header_len
                
                # Mobi típus
                mobi_type = struct.unpack('>I', data[mobi_header_start+8:mobi_header_start+12])[0]
                info["mobi_type"] = mobi_type
                
                # Encoding
                encoding = struct.unpack('>I', data[mobi_header_start+28:mobi_header_start+32])[0]
                info["encoding"] = "UTF-8" if encoding == 65001 else f"Code {encoding}"
                
            except Exception as e:
                info["metadata_error"] = str(e)
            
            return info
            
    except Exception as e:
        return {"error": f"Cannot analyze MOBI structure: {str(e)}"}
